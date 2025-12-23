"""
Core PowerPoint generation functions.
"""

import os
import tempfile
import shutil
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

from . import config_loader
config = config_loader.config


def hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def copy_to_temp(file_path):
    """
    Copy file to temporary location to avoid OneDrive locking issues.
    
    Args:
        file_path: Path to source file
    
    Returns:
        Path to temporary file, or original path if copy fails
    """
    try:
        temp_dir = tempfile.gettempdir()
        # Create unique temp filename to avoid conflicts
        base_name = os.path.basename(file_path)
        name, ext = os.path.splitext(base_name)
        temp_file = tempfile.NamedTemporaryFile(
            mode='wb',
            suffix=ext,
            prefix=f"{name}_",
            delete=False,
            dir=temp_dir
        )
        temp_file_path = temp_file.name
        temp_file.close()
        
        # Copy file to temp location
        shutil.copy2(file_path, temp_file_path)
        return temp_file_path
    except Exception as e:
        print(f"Warning: Could not copy file to temp location: {e}")
        print("Attempting to read from original location...")
        return file_path


def read_objectives(excel_file):
    """
    Read Objectives sheet from Excel file.
    Returns: DataFrame with 'North star' and 'Key elements' columns
    """
    # Copy to temp location to avoid OneDrive locking issues
    temp_file = copy_to_temp(excel_file)
    try:
        df = pd.read_excel(temp_file, sheet_name='Objectives')
        # Handle case-insensitive column names
        df.columns = df.columns.str.strip()
        # Try to find columns (case-insensitive)
        north_star_col = None
        key_elements_col = None
        
        for col in df.columns:
            col_lower = col.lower()
            if 'north' in col_lower and 'star' in col_lower:
                north_star_col = col
            elif 'key' in col_lower and 'element' in col_lower:
                key_elements_col = col
        
        if north_star_col is None:
            # Assume first column is North star
            north_star_col = df.columns[0]
        if key_elements_col is None:
            # Assume second column is Key elements
            key_elements_col = df.columns[1] if len(df.columns) > 1 else None
        
        # Filter out empty rows
        df = df.dropna(subset=[north_star_col])
        
        return {
            'north_star': df[north_star_col].iloc[0] if len(df) > 0 else "",
            'key_elements': df[key_elements_col].dropna().tolist() if key_elements_col else []
        }
    except Exception as e:
        print(f"Error reading Objectives sheet: {e}")
        return {'north_star': '', 'key_elements': []}
    finally:
        # Clean up temp file if it was created
        if temp_file != excel_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except Exception as e:
                print(f"Warning: Could not delete temp file {temp_file}: {e}")


def read_roadmap(excel_file):
    """
    Read Roadmap sheet from Excel file.
    Returns: DataFrame with Timeline, Phase, and Workpackage columns
    """
    # Copy to temp location to avoid OneDrive locking issues
    temp_file = copy_to_temp(excel_file)
    try:
        df = pd.read_excel(temp_file, sheet_name='Roadmap')
        df.columns = df.columns.str.strip()
        
        # Try to identify columns (case-insensitive)
        timeline_col = None
        phase_col = None
        workpackage_col = None
        
        for col in df.columns:
            col_lower = col.lower()
            if 'timeline' in col_lower or 'phase' in col_lower and timeline_col is None:
                timeline_col = col
            elif 'phase' in col_lower and phase_col is None:
                phase_col = col
            elif 'workpackage' in col_lower or 'work package' in col_lower:
                workpackage_col = col
        
        # Fallback to positional columns
        if timeline_col is None:
            timeline_col = df.columns[0]
        if phase_col is None:
            phase_col = df.columns[1] if len(df.columns) > 1 else None
        if workpackage_col is None:
            workpackage_col = df.columns[2] if len(df.columns) > 2 else None
        
        # Filter out empty rows
        df = df.dropna(subset=[timeline_col])
        
        result_df = pd.DataFrame({
            'Timeline': df[timeline_col],
            'Phase': df[phase_col] if phase_col else [''] * len(df),
            'Workpackage': df[workpackage_col] if workpackage_col else [''] * len(df)
        })
        
        return result_df
    except Exception as e:
        print(f"Error reading Roadmap sheet: {e}")
        return pd.DataFrame(columns=['Timeline', 'Phase', 'Workpackage'])
    finally:
        # Clean up temp file if it was created
        if temp_file != excel_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except Exception as e:
                print(f"Warning: Could not delete temp file {temp_file}: {e}")


def add_logo(slide, logo_path, position):
    """Add logo to slide at specified position."""
    if logo_path is None or not os.path.exists(logo_path):
        return
    
    try:
        left_map = {
            'top_left': Inches(0.5),
            'top_right': config.SLIDE_WIDTH - Inches(1.5),
            'bottom_left': Inches(0.5),
            'bottom_right': config.SLIDE_WIDTH - Inches(1.5),
            'center': (config.SLIDE_WIDTH - Inches(1)) / 2
        }
        top_map = {
            'top_left': Inches(0.3),
            'top_right': Inches(0.3),
            'bottom_left': config.SLIDE_HEIGHT - Inches(0.8),
            'bottom_right': config.SLIDE_HEIGHT - Inches(0.8),
            'center': Inches(0.3)
        }
        
        left = left_map.get(position, Inches(0.5))
        top = top_map.get(position, Inches(0.3))
        
        slide.shapes.add_picture(logo_path, left, top, height=Inches(0.7))
    except Exception as e:
        print(f"Warning: Could not add logo: {e}")


def load_template_slide(template_path, slide_index=0):
    """
    Load a slide from a template PowerPoint file (.pptx or .potx).
    
    Args:
        template_path: Path to template file
        slide_index: Index of slide to use from template (default: 0)
    
    Returns:
        Presentation object with template slide, or None if loading fails
    """
    if template_path is None:
        return None
    
    if not os.path.exists(template_path):
        print(f"Warning: Template file not found: {template_path}")
        return None
    
    try:
        template_prs = Presentation(template_path)
        if slide_index >= len(template_prs.slides):
            print(f"Warning: Slide index {slide_index} out of range for template {template_path} (has {len(template_prs.slides)} slides)")
            return None
        return template_prs
    except Exception as e:
        print(f"Warning: Could not load template file {template_path}: {e}")
        return None


def create_slide_from_template(prs, template_prs, slide_index=0):
    """
    Create a new slide in the presentation using a template slide.
    
    Args:
        prs: Target presentation to add slide to
        template_prs: Template presentation object
        slide_index: Index of slide from template to copy
    
    Returns:
        New slide object, or None if template is invalid
    """
    if template_prs is None:
        return None
    
    try:
        if slide_index >= len(template_prs.slides):
            return None
        
        template_slide = template_prs.slides[slide_index]
        
        # Use blank layout from target presentation (layouts can't be reused across presentations)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Copy background from template first
        try:
            if hasattr(template_slide.background, 'fill'):
                if template_slide.background.fill.type == 1:  # Solid fill
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = template_slide.background.fill.fore_color.rgb
        except (AttributeError, ValueError):
            pass
        
        # Copy all shapes from template slide
        for shape in template_slide.shapes:
            # Get shape properties
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Copy shape based on type
            if hasattr(shape, 'image'):
                # Picture shape - wrap blob in BytesIO
                try:
                    image = shape.image
                    image_blob = BytesIO(image.blob)
                    slide.shapes.add_picture(image_blob, left, top, width, height)
                except (AttributeError, ValueError, IOError):
                    pass
            elif hasattr(shape, 'text'):
                # Text shape - copy paragraphs individually to preserve structure
                try:
                    new_shape = slide.shapes.add_textbox(left, top, width, height)
                    new_text_frame = new_shape.text_frame
                    new_text_frame.word_wrap = shape.text_frame.word_wrap
                    
                    # Clear default paragraph
                    new_text_frame.clear()
                    
                    # Copy each paragraph from template
                    for para in shape.text_frame.paragraphs:
                        # Add new paragraph
                        new_para = new_text_frame.add_paragraph()
                        
                        # Copy paragraph-level properties
                        new_para.level = para.level
                        new_para.alignment = para.alignment
                        new_para.space_before = para.space_before
                        new_para.space_after = para.space_after
                        
                        # Copy runs within paragraph (preserves formatting)
                        for run in para.runs:
                            new_run = new_para.add_run()
                            new_run.text = run.text
                            # Copy run formatting
                            new_run.font.name = run.font.name
                            new_run.font.size = run.font.size
                            new_run.font.bold = run.font.bold
                            new_run.font.italic = run.font.italic
                            new_run.font.underline = run.font.underline
                            try:
                                new_run.font.color.rgb = run.font.color.rgb
                            except (AttributeError, ValueError):
                                pass
                        
                        # If no runs, copy paragraph text directly
                        if len(para.runs) == 0:
                            new_para.text = para.text
                            # Apply paragraph-level font formatting
                            try:
                                new_para.font.name = para.font.name
                                new_para.font.size = para.font.size
                                new_para.font.bold = para.font.bold
                                new_para.font.color.rgb = para.font.color.rgb
                            except (AttributeError, ValueError):
                                pass
                except (AttributeError, ValueError, IndexError):
                    # Fallback: simple text copy if detailed copying fails
                    try:
                        new_shape = slide.shapes.add_textbox(left, top, width, height)
                        new_shape.text_frame.text = shape.text_frame.text
                    except (AttributeError, ValueError):
                        pass
        
        return slide
    except Exception as e:
        print(f"Warning: Could not copy template slide: {e}")
        return None


def calculate_text_height(text, width, font_size, min_height=None, max_height=None):
    """
    Calculate estimated height needed for text based on content length and width.
    
    Args:
        text: Text content
        width: Available width in inches
        font_size: Font size in points
        min_height: Minimum height in inches (optional)
        max_height: Maximum height in inches (optional)
    
    Returns:
        Estimated height in inches
    """
    if not text:
        return min_height or Inches(0.5)
    
    # Estimate characters per line based on width and font size
    # Rough estimate: 1 inch ≈ 12 characters at 18pt font
    chars_per_inch = max(8, 12 * (18 / font_size.pt))
    chars_per_line = int(width / Inches(1) * chars_per_inch)
    
    # Calculate estimated lines
    text_length = len(str(text))
    estimated_lines = max(1, (text_length // chars_per_line) + 1)
    
    # Add some buffer for word wrapping
    estimated_lines = int(estimated_lines * 1.2)
    
    # Calculate height (roughly 1.2x font size per line)
    line_height = font_size.pt * 1.2 / 72  # Convert points to inches
    estimated_height = Inches(estimated_lines * line_height)
    
    # Apply min/max constraints
    if min_height and estimated_height < min_height:
        estimated_height = min_height
    if max_height and estimated_height > max_height:
        estimated_height = max_height
    
    return estimated_height


def create_title_slide(prs, objectives_data, title="Roadmap Presentation"):
    """Create title slide with branding."""
    # Try to use template if configured
    template_prs = load_template_slide(config.TITLE_SLIDE_TEMPLATE, config.TEMPLATE_SLIDE_INDEX)
    if template_prs:
        slide = create_slide_from_template(prs, template_prs, config.TEMPLATE_SLIDE_INDEX)
        if slide is None:
            # Fall back to blank layout if template copy failed
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
    else:
        # Use blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
    
    # Add logo if configured (only if not using template or template doesn't have logo)
    if config.LOGO_PATH:
        add_logo(slide, config.LOGO_PATH, config.LOGO_POSITION)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        config.SIDE_MARGIN,
        config.TITLE_TOP_MARGIN,
        config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
        Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = config.TITLE_FONT_NAME
    title_para.font.size = config.TITLE_FONT_SIZE
    title_para.font.bold = True
    title_para.font.color.rgb = config.BRAND_PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle (North star)
    if objectives_data.get('north_star'):
        subtitle_box = slide.shapes.add_textbox(
            config.SIDE_MARGIN,
            config.TITLE_TOP_MARGIN + Inches(2.5),
            config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
            config.SLIDE_HEIGHT - config.TITLE_TOP_MARGIN - Inches(2.5) - config.BOTTOM_MARGIN
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = objectives_data['north_star']
        subtitle_frame.word_wrap = True
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = config.BODY_FONT_NAME
        subtitle_para.font.size = config.SUBTITLE_FONT_SIZE
        subtitle_para.font.color.rgb = config.BRAND_SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER


def create_objectives_slide(prs, objectives_data):
    """Create objectives slide(s) displaying North star and key elements with pagination."""
    key_elements = objectives_data.get('key_elements', [])
    has_north_star = bool(objectives_data.get('north_star'))
    
    # Calculate available space for key elements
    # Account for: slide title, north star section (if present), key elements title
    SLIDE_TITLE_HEIGHT = Inches(0.8)
    SLIDE_TITLE_TOP = Inches(0.5)
    KEY_ELEMENTS_TITLE_HEIGHT = Inches(0.6)
    KEY_ELEMENTS_TITLE_SPACING = Inches(0.8)
    
    # Calculate space after North Star (if present)
    if has_north_star:
        north_star_text = str(objectives_data['north_star'])
        box_width = config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN) - (2 * config.TEXT_BOX_MARGIN)
        north_star_height = calculate_text_height(
            north_star_text,
            box_width,
            config.BODY_FONT_SIZE,
            min_height=config.NORTH_STAR_MIN_HEIGHT,
            max_height=config.NORTH_STAR_MAX_HEIGHT
        )
        space_after_north_star = Inches(1.2) + north_star_height + Inches(0.3)  # Header + content + spacing
    else:
        space_after_north_star = 0
    
    # Calculate available height for key elements
    available_height = (
        config.SLIDE_HEIGHT 
        - SLIDE_TITLE_TOP 
        - SLIDE_TITLE_HEIGHT 
        - space_after_north_star
        - KEY_ELEMENTS_TITLE_HEIGHT 
        - KEY_ELEMENTS_TITLE_SPACING
        - config.BOTTOM_MARGIN
    )
    
    # Calculate how many key elements fit per slide
    items_per_slide = max(1, int(available_height / config.KEY_ELEMENT_HEIGHT_ESTIMATE))
    total_elements = len(key_elements)
    slides_needed = max(1, (total_elements + items_per_slide - 1) // items_per_slide) if total_elements > 0 else 1
    
    # Create slides
    # Try to use template if configured
    template_prs = load_template_slide(config.CONTENT_SLIDE_TEMPLATE, config.TEMPLATE_SLIDE_INDEX)
    
    for slide_num in range(slides_needed):
        # Use template if available, otherwise use blank layout
        if template_prs:
            slide = create_slide_from_template(prs, template_prs, config.TEMPLATE_SLIDE_INDEX)
            if slide is None:
                # Fall back to blank layout if template copy failed
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Set background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
        else:
            # Use blank layout
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
        
        # Add logo (only if not using template or template doesn't have logo)
        if config.LOGO_PATH:
            add_logo(slide, config.LOGO_PATH, config.LOGO_POSITION)
        
        # Slide title - add page number if multiple slides
        title_text = "Objectives"
        if slides_needed > 1:
            title_text += f" (Page {slide_num + 1} of {slides_needed})"
        
        title_box = slide.shapes.add_textbox(
            config.SIDE_MARGIN,
            SLIDE_TITLE_TOP,
            config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
            SLIDE_TITLE_HEIGHT
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.name = config.TITLE_FONT_NAME
        title_para.font.size = config.HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = config.BRAND_PRIMARY_COLOR
        
        y_pos = config.CONTENT_TOP_MARGIN
        
        # North Star section - only show on first slide
        if has_north_star and slide_num == 0:
            north_star_box = slide.shapes.add_textbox(
                config.SIDE_MARGIN,
                y_pos,
                config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
                Inches(1)
            )
            north_star_frame = north_star_box.text_frame
            north_star_frame.text = "North Star"
            north_star_para = north_star_frame.paragraphs[0]
            north_star_para.font.name = config.BODY_FONT_NAME
            north_star_para.font.size = Pt(24)
            north_star_para.font.bold = True
            north_star_para.font.color.rgb = config.BRAND_SECONDARY_COLOR
            
            y_pos += Inches(1.2)
            
            # Calculate dynamic height for North Star content box
            north_star_text = str(objectives_data['north_star'])
            box_width = config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN) - (2 * config.TEXT_BOX_MARGIN)
            dynamic_height = calculate_text_height(
                north_star_text,
                box_width,
                config.BODY_FONT_SIZE,
                min_height=config.NORTH_STAR_MIN_HEIGHT,
                max_height=config.NORTH_STAR_MAX_HEIGHT
            )
            
            # North star content box
            if config.USE_SHAPES:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    config.SIDE_MARGIN,
                    y_pos,
                    config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
                    dynamic_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = config.CONTENT_BOX_COLOR
                shape.line.color.rgb = config.BRAND_SECONDARY_COLOR
                shape.line.width = Pt(2)
                
                text_frame = shape.text_frame
                text_frame.text = north_star_text
                text_frame.word_wrap = True
                text_frame.margin_left = config.TEXT_BOX_MARGIN
                text_frame.margin_right = config.TEXT_BOX_MARGIN
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                para = text_frame.paragraphs[0]
                para.font.name = config.BODY_FONT_NAME
                para.font.size = config.BODY_FONT_SIZE
                para.font.color.rgb = config.BRAND_TEXT_COLOR
                para.alignment = PP_ALIGN.LEFT
            else:
                north_star_content = slide.shapes.add_textbox(
                    config.SIDE_MARGIN,
                    y_pos,
                    config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
                    dynamic_height
                )
                north_star_content_frame = north_star_content.text_frame
                north_star_content_frame.text = north_star_text
                north_star_content_frame.word_wrap = True
                north_star_content_para = north_star_content_frame.paragraphs[0]
                north_star_content_para.font.name = config.BODY_FONT_NAME
                north_star_content_para.font.size = config.BODY_FONT_SIZE
                north_star_content_para.font.color.rgb = config.BRAND_TEXT_COLOR
                north_star_content_para.alignment = PP_ALIGN.LEFT
            
            y_pos += dynamic_height + Inches(0.3)
        
        # Key Elements section
        if total_elements > 0:
            key_elements_title = slide.shapes.add_textbox(
                config.SIDE_MARGIN,
                y_pos,
                config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
                KEY_ELEMENTS_TITLE_HEIGHT
            )
            key_elements_title_frame = key_elements_title.text_frame
            key_elements_title_frame.text = "Key Elements"
            key_elements_title_para = key_elements_title_frame.paragraphs[0]
            key_elements_title_para.font.name = config.BODY_FONT_NAME
            key_elements_title_para.font.size = Pt(24)
            key_elements_title_para.font.bold = True
            key_elements_title_para.font.color.rgb = config.BRAND_SECONDARY_COLOR
            
            y_pos += KEY_ELEMENTS_TITLE_SPACING
            
            # Calculate which elements to show on this slide
            start_idx = slide_num * items_per_slide
            end_idx = min(start_idx + items_per_slide, total_elements)
            elements_for_slide = key_elements[start_idx:end_idx]
            
            # Calculate available height for elements list
            elements_height = config.SLIDE_HEIGHT - y_pos - config.BOTTOM_MARGIN
            
            # Key elements list
            elements_box = slide.shapes.add_textbox(
                config.SIDE_MARGIN + Inches(0.3),
                y_pos,
                config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN) - Inches(0.3),
                elements_height
            )
            elements_frame = elements_box.text_frame
            elements_frame.word_wrap = True
            
            for i, element in enumerate(elements_for_slide):
                if i > 0:
                    elements_frame.add_paragraph()
                para = elements_frame.paragraphs[i]
                para.text = f"• {str(element)}"
                para.font.name = config.BODY_FONT_NAME
                para.font.size = config.BODY_FONT_SIZE
                para.font.color.rgb = config.BRAND_TEXT_COLOR
                para.space_after = Pt(8)
                para.level = 0
                para.alignment = PP_ALIGN.LEFT


def create_timeline_overview_slide(prs, roadmap_df):
    """Create timeline overview slide showing horizontal flow of timeline steps and phases."""
    if roadmap_df.empty:
        return
    
    # Try to use template if configured
    template_prs = load_template_slide(config.CONTENT_SLIDE_TEMPLATE, config.TEMPLATE_SLIDE_INDEX)
    
    if template_prs:
        slide = create_slide_from_template(prs, template_prs, config.TEMPLATE_SLIDE_INDEX)
        if slide is None:
            # Fall back to blank layout if template copy failed
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
    else:
        # Use blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
    
    # Add logo
    if config.LOGO_PATH:
        add_logo(slide, config.LOGO_PATH, config.LOGO_POSITION)
    
    # Slide title
    title_box = slide.shapes.add_textbox(
        config.SIDE_MARGIN,
        Inches(0.5),
        config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Roadmap Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = config.TITLE_FONT_NAME
    title_para.font.size = config.HEADING_FONT_SIZE
    title_para.font.bold = True
    title_para.font.color.rgb = config.BRAND_PRIMARY_COLOR
    
    # Group by Timeline, then collect phases for each timeline
    timeline_groups = roadmap_df.groupby('Timeline', sort=False)
    
    # Collect timeline-phase items in order
    timeline_phase_items = []
    for timeline, group in timeline_groups:
        phases = group['Phase'].dropna().unique().tolist()
        if len(phases) > 0:
            for phase in phases:
                if phase and str(phase).strip():
                    timeline_phase_items.append({
                        'timeline': str(timeline),
                        'phase': str(phase)
                    })
        else:
            # Timeline with no phases
            timeline_phase_items.append({
                'timeline': str(timeline),
                'phase': None
            })
    
    if not timeline_phase_items:
        return
    
    # Calculate layout dimensions using config (increased size)
    shape_width = config.OVERVIEW_SHAPE_WIDTH * 1.25  # Make 25% larger
    shape_height = config.OVERVIEW_SHAPE_HEIGHT * 1.25  # Make 25% larger
    overlap_amount = Inches(0.6)  # Amount pentagons overlap to fit into each other
    
    # Calculate total width needed (accounting for overlap)
    total_items = len(timeline_phase_items)
    total_width = (total_items * shape_width) - (overlap_amount * (total_items - 1))
    
    # Start position (centered)
    available_width = config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN)
    start_x = config.SIDE_MARGIN + (available_width - total_width) / 2 if total_width < available_width else config.SIDE_MARGIN
    y_pos = config.CONTENT_TOP_MARGIN
    
    current_x = start_x
    
    # Draw timeline shapes fitting into each other
    for i, item in enumerate(timeline_phase_items):
        timeline = item['timeline']
        phase = item['phase']
        
        # Create text content: Timeline with phase as subtext
        if phase:
            box_text = f"{timeline}\n{phase}"
        else:
            box_text = timeline
        
        # Timeline shape using PENTAGON AutoShape
        timeline_shape = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON,
            current_x,
            y_pos,
            shape_width,
            shape_height
        )
        timeline_shape.fill.solid()
        timeline_shape.fill.fore_color.rgb = config.OVERVIEW_TIMELINE_SHAPE_COLOR
        timeline_shape.line.color.rgb = config.OVERVIEW_TIMELINE_SHAPE_COLOR
        timeline_shape.line.width = Pt(2)
        
        timeline_text = timeline_shape.text_frame
        timeline_text.text = box_text
        timeline_text.word_wrap = True
        timeline_text.margin_left = Inches(0.1)
        timeline_text.margin_right = Inches(0.1)
        timeline_text.margin_top = Inches(0.1)
        timeline_text.margin_bottom = Inches(0.1)
        
        # Timeline line (bold)
        timeline_para = timeline_text.paragraphs[0]
        timeline_para.font.name = config.BODY_FONT_NAME
        timeline_para.font.size = Pt(18)
        timeline_para.font.bold = True
        timeline_para.font.color.rgb = config.OVERVIEW_TIMELINE_TEXT_COLOR
        timeline_para.alignment = PP_ALIGN.CENTER
        
        # Phase subtext line (if exists)
        if phase:
            if len(timeline_text.paragraphs) < 2:
                timeline_text.add_paragraph()
            phase_para = timeline_text.paragraphs[1]
            phase_para.text = phase
            phase_para.font.name = config.BODY_FONT_NAME
            phase_para.font.size = Pt(12)
            phase_para.font.bold = False
            phase_para.font.color.rgb = config.OVERVIEW_TIMELINE_TEXT_COLOR
            phase_para.alignment = PP_ALIGN.CENTER
        
        # Move to next position with overlap (fit into next shape)
        current_x += shape_width - overlap_amount


def create_roadmap_slides(prs, roadmap_df):
    """Create roadmap slides grouped by timeline/phase with pagination support."""
    if roadmap_df.empty:
        return
    
    # Constants for pagination calculations
    SLIDE_TITLE_HEIGHT = Inches(0.8)
    SLIDE_TITLE_TOP = Inches(0.5)
    PHASE_HEADER_HEIGHT = Inches(0.7)
    ITEM_HEIGHT_ESTIMATE = Inches(0.5)  # Estimated height per workpackage item
    MIN_ITEM_HEIGHT = Inches(0.3)
    MAX_CONTENT_HEIGHT = config.SLIDE_HEIGHT - config.CONTENT_TOP_MARGIN - config.BOTTOM_MARGIN
    
    # Group by Timeline
    grouped = roadmap_df.groupby('Timeline')
    
    for timeline, group in grouped:
        # Group by Phase within this timeline
        phase_groups = list(group.groupby('Phase') if 'Phase' in group.columns else [(None, group)])
        num_phases = len(phase_groups)
        
        # Calculate layout dimensions
        max_width = (config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN)) / num_phases if num_phases > 1 else (config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN))
        box_width = max_width - Inches(0.2)
        
        # Organize workpackages by phase with pagination info
        phase_data_list = []
        for phase, phase_data in phase_groups:
            workpackages = phase_data['Workpackage'].dropna().tolist()
            phase_data_list.append({
                'phase': phase,
                'workpackages': workpackages,
                'total_items': len(workpackages)
            })
        
        # Calculate how many items can fit per slide
        available_height = MAX_CONTENT_HEIGHT - PHASE_HEADER_HEIGHT
        items_per_slide = max(1, int(available_height / ITEM_HEIGHT_ESTIMATE))
        
        # Determine if we need multiple slides
        max_items_any_phase = max([pd['total_items'] for pd in phase_data_list], default=0)
        slides_needed = max(1, (max_items_any_phase + items_per_slide - 1) // items_per_slide)
        
        # Try to use template if configured
        template_prs = load_template_slide(config.CONTENT_SLIDE_TEMPLATE, config.TEMPLATE_SLIDE_INDEX)
        
        # Create slides for this timeline
        for slide_num in range(slides_needed):
            # Use template if available, otherwise use blank layout
            if template_prs:
                slide = create_slide_from_template(prs, template_prs, config.TEMPLATE_SLIDE_INDEX)
                if slide is None:
                    # Fall back to blank layout if template copy failed
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # Set background
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
            else:
                # Use blank layout
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Set background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = config.BRAND_BACKGROUND_COLOR
            
            # Add logo (only if not using template or template doesn't have logo)
            if config.LOGO_PATH:
                add_logo(slide, config.LOGO_PATH, config.LOGO_POSITION)
            
            # Slide title (Timeline) - add page number if multiple slides
            title_text = f"Roadmap: {str(timeline)}"
            if slides_needed > 1:
                title_text += f" (Page {slide_num + 1} of {slides_needed})"
            
            title_box = slide.shapes.add_textbox(
                config.SIDE_MARGIN,
                SLIDE_TITLE_TOP,
                config.SLIDE_WIDTH - (2 * config.SIDE_MARGIN),
                SLIDE_TITLE_HEIGHT
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_para = title_frame.paragraphs[0]
            title_para.font.name = config.TITLE_FONT_NAME
            title_para.font.size = config.HEADING_FONT_SIZE
            title_para.font.bold = True
            title_para.font.color.rgb = config.BRAND_PRIMARY_COLOR
            
            y_pos = config.CONTENT_TOP_MARGIN
            
            # Add content for each phase
            for phase_idx, phase_info in enumerate(phase_data_list):
                phase = phase_info['phase']
                workpackages = phase_info['workpackages']
                
                # Calculate which items to show on this slide
                start_idx = slide_num * items_per_slide
                end_idx = min(start_idx + items_per_slide, len(workpackages))
                
                items_for_this_slide = workpackages[start_idx:end_idx] if start_idx < len(workpackages) else []
                
                x_pos = config.SIDE_MARGIN + (phase_idx * max_width) + Inches(0.1)
                
                # Phase header - always show if phase name exists
                if phase and str(phase).strip():
                    phase_title = slide.shapes.add_textbox(
                        x_pos,
                        y_pos,
                        box_width,
                        Inches(0.6)
                    )
                    phase_title_frame = phase_title.text_frame
                    phase_title_frame.text = str(phase)
                    phase_title_frame.word_wrap = True  # Ensure phase titles wrap if too long
                    phase_title_para = phase_title_frame.paragraphs[0]
                    phase_title_para.font.name = config.BODY_FONT_NAME
                    phase_title_para.font.size = Pt(22)
                    phase_title_para.font.bold = True
                    phase_title_para.font.color.rgb = config.BRAND_SECONDARY_COLOR
                    phase_title_para.alignment = PP_ALIGN.LEFT
                    y_pos_phase = y_pos + PHASE_HEADER_HEIGHT
                else:
                    y_pos_phase = y_pos
                
                # Only show content box if there are items for this slide
                if items_for_this_slide:
                    # Calculate content height based on number of items
                    content_height = min(
                        MAX_CONTENT_HEIGHT - (y_pos_phase - config.CONTENT_TOP_MARGIN),
                        max(MIN_ITEM_HEIGHT, len(items_for_this_slide) * ITEM_HEIGHT_ESTIMATE + Inches(0.3))
                    )
                    if config.USE_SHAPES:
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            x_pos,
                            y_pos_phase,
                            box_width,
                            content_height
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = config.CONTENT_BOX_COLOR
                        shape.line.color.rgb = config.BRAND_ACCENT_COLOR
                        shape.line.width = Pt(1.5)
                        
                        text_frame = shape.text_frame
                        text_frame.word_wrap = True
                        text_frame.margin_left = Inches(0.15)
                        text_frame.margin_right = Inches(0.15)
                        text_frame.margin_top = Inches(0.15)
                        text_frame.margin_bottom = Inches(0.15)
                        
                        for i, wp in enumerate(items_for_this_slide):
                            if i > 0:
                                text_frame.add_paragraph()
                            para = text_frame.paragraphs[i]
                            para.text = f"• {str(wp)}"
                            para.font.name = config.BODY_FONT_NAME
                            para.font.size = Pt(14)
                            para.font.color.rgb = config.BRAND_TEXT_COLOR
                            para.space_after = Pt(6)
                            para.alignment = PP_ALIGN.LEFT
                    else:
                        wp_box = slide.shapes.add_textbox(
                            x_pos,
                            y_pos_phase,
                            box_width,
                            content_height
                        )
                        wp_frame = wp_box.text_frame
                        wp_frame.word_wrap = True
                        
                        for i, wp in enumerate(items_for_this_slide):
                            if i > 0:
                                wp_frame.add_paragraph()
                            para = wp_frame.paragraphs[i]
                            para.text = f"• {str(wp)}"
                            para.font.name = config.BODY_FONT_NAME
                            para.font.size = Pt(14)
                            para.font.color.rgb = config.BRAND_TEXT_COLOR
                            para.space_after = Pt(6)
                            para.alignment = PP_ALIGN.LEFT


def generate_presentation(excel_file, output_path=None):
    """
    Generate PowerPoint presentation from Excel file.
    
    Args:
        excel_file: Path to Excel file
        output_path: Optional output path (defaults to same name as Excel with .pptx extension)
    
    Returns:
        Path to generated PowerPoint file
    """
    import os
    
    # Determine output file path
    if output_path is None:
        base_name = os.path.splitext(excel_file)[0]
        output_path = f"{base_name}.pptx"
    
    # Extract Excel filename (without extension) for title
    excel_filename = os.path.splitext(os.path.basename(excel_file))[0]
    
    # Read Excel data
    objectives_data = read_objectives(excel_file)
    roadmap_df = read_roadmap(excel_file)
    
    print(f"Found {len(roadmap_df)} roadmap entries")
    print(f"North Star: {objectives_data.get('north_star', 'Not found')}")
    print(f"Key Elements: {len(objectives_data.get('key_elements', []))} items")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = config.SLIDE_WIDTH
    prs.slide_height = config.SLIDE_HEIGHT
    
    # Generate slides
    print("Generating slides...")
    create_title_slide(prs, objectives_data, title=excel_filename)
    create_objectives_slide(prs, objectives_data)
    create_timeline_overview_slide(prs, roadmap_df)
    create_roadmap_slides(prs, roadmap_df)
    
    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint presentation saved to: {output_path}")
    
    return output_path

