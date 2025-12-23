"""
Excel to PowerPoint Roadmap Generator
Reads Excel file with Objectives and Roadmap sheets and generates a branded PowerPoint presentation.
"""

import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd


################################################################################
### CONFIGURATION SECTION - CUSTOMIZE YOUR BRANDING HERE ###
################################################################################

# BRANDING CONFIG - Colors (RGB values 0-255)
BRAND_PRIMARY_COLOR = RGBColor(0, 51, 102)      # Deep blue - CHANGE THIS
BRAND_SECONDARY_COLOR = RGBColor(0, 102, 204)   # Medium blue - CHANGE THIS
BRAND_ACCENT_COLOR = RGBColor(255, 153, 0)      # Orange accent - CHANGE THIS
BRAND_TEXT_COLOR = RGBColor(51, 51, 51)         # Dark gray text - CHANGE THIS
BRAND_BACKGROUND_COLOR = RGBColor(255, 255, 255) # White background - CHANGE THIS

# Logo Configuration
LOGO_PATH = None  # Set to path of your logo file (e.g., "logo.png") or None to skip
LOGO_POSITION = "top_right"  # Options: "top_left", "top_right", "bottom_left", "bottom_right", "center"

# Font Configuration
TITLE_FONT_NAME = "Calibri"  # CHANGE THIS to your brand font
BODY_FONT_NAME = "Calibri"   # CHANGE THIS to your brand font
TITLE_FONT_SIZE = Pt(44)
SUBTITLE_FONT_SIZE = Pt(28)
HEADING_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(18)

# SLIDE LAYOUT CONFIG - Dimensions and spacing
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
TITLE_TOP_MARGIN = Inches(1)
CONTENT_TOP_MARGIN = Inches(1.5)
SIDE_MARGIN = Inches(0.5)
BOTTOM_MARGIN = Inches(0.5)

# VISUAL STYLE CONFIG
USE_SHAPES = True  # Use rounded rectangles for content boxes
SHAPE_CORNER_RADIUS = Inches(0.1)
CONTENT_BOX_COLOR = RGBColor(245, 245, 245)  # Light gray for content boxes

################################################################################
### END CONFIGURATION SECTION ###
################################################################################


def hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def read_objectives(excel_file):
    """
    Read Objectives sheet from Excel file.
    Returns: DataFrame with 'North star' and 'Key elements' columns
    """
    try:
        df = pd.read_excel(excel_file, sheet_name='Objectives')
        # Handle case-insensitive column names
        df.columns = df.columns.str.strip()
        # Try to find columns (case-insensitive)
        north_star_col = None
        key_elements_col = None
        
        for col in df.columns:
            col_lower = col.lower()
            if 'north' in col_lower and 'star' in col_lower:
                north_star_col = col
            elif 'key' in col_lower and 'element' in col_lower:
                key_elements_col = col
        
        if north_star_col is None:
            # Assume first column is North star
            north_star_col = df.columns[0]
        if key_elements_col is None:
            # Assume second column is Key elements
            key_elements_col = df.columns[1] if len(df.columns) > 1 else None
        
        # Filter out empty rows
        df = df.dropna(subset=[north_star_col])
        
        return {
            'north_star': df[north_star_col].iloc[0] if len(df) > 0 else "",
            'key_elements': df[key_elements_col].dropna().tolist() if key_elements_col else []
        }
    except Exception as e:
        print(f"Error reading Objectives sheet: {e}")
        return {'north_star': '', 'key_elements': []}


def read_roadmap(excel_file):
    """
    Read Roadmap sheet from Excel file.
    Returns: DataFrame with Timeline, Phase, and Workpackage columns
    """
    try:
        df = pd.read_excel(excel_file, sheet_name='Roadmap')
        df.columns = df.columns.str.strip()
        
        # Try to identify columns (case-insensitive)
        timeline_col = None
        phase_col = None
        workpackage_col = None
        
        for col in df.columns:
            col_lower = col.lower()
            if 'timeline' in col_lower or 'phase' in col_lower and timeline_col is None:
                timeline_col = col
            elif 'phase' in col_lower and phase_col is None:
                phase_col = col
            elif 'workpackage' in col_lower or 'work package' in col_lower:
                workpackage_col = col
        
        # Fallback to positional columns
        if timeline_col is None:
            timeline_col = df.columns[0]
        if phase_col is None:
            phase_col = df.columns[1] if len(df.columns) > 1 else None
        if workpackage_col is None:
            workpackage_col = df.columns[2] if len(df.columns) > 2 else None
        
        # Filter out empty rows
        df = df.dropna(subset=[timeline_col])
        
        result_df = pd.DataFrame({
            'Timeline': df[timeline_col],
            'Phase': df[phase_col] if phase_col else [''] * len(df),
            'Workpackage': df[workpackage_col] if workpackage_col else [''] * len(df)
        })
        
        return result_df
    except Exception as e:
        print(f"Error reading Roadmap sheet: {e}")
        return pd.DataFrame(columns=['Timeline', 'Phase', 'Workpackage'])


def add_logo(slide, logo_path, position):
    """Add logo to slide at specified position."""
    if logo_path is None or not os.path.exists(logo_path):
        return
    
    try:
        left_map = {
            'top_left': Inches(0.5),
            'top_right': SLIDE_WIDTH - Inches(1.5),
            'bottom_left': Inches(0.5),
            'bottom_right': SLIDE_WIDTH - Inches(1.5),
            'center': (SLIDE_WIDTH - Inches(1)) / 2
        }
        top_map = {
            'top_left': Inches(0.3),
            'top_right': Inches(0.3),
            'bottom_left': SLIDE_HEIGHT - Inches(0.8),
            'bottom_right': SLIDE_HEIGHT - Inches(0.8),
            'center': Inches(0.3)
        }
        
        left = left_map.get(position, Inches(0.5))
        top = top_map.get(position, Inches(0.3))
        
        slide.shapes.add_picture(logo_path, left, top, height=Inches(0.7))
    except Exception as e:
        print(f"Warning: Could not add logo: {e}")


def create_title_slide(prs, objectives_data):
    """Create title slide with branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_BACKGROUND_COLOR
    
    # Add logo if configured
    if LOGO_PATH:
        add_logo(slide, LOGO_PATH, LOGO_POSITION)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        SIDE_MARGIN,
        TITLE_TOP_MARGIN,
        SLIDE_WIDTH - (2 * SIDE_MARGIN),
        Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Roadmap Presentation"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = TITLE_FONT_NAME
    title_para.font.size = TITLE_FONT_SIZE
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle (North star)
    if objectives_data.get('north_star'):
        subtitle_box = slide.shapes.add_textbox(
            SIDE_MARGIN,
            TITLE_TOP_MARGIN + Inches(2.5),
            SLIDE_WIDTH - (2 * SIDE_MARGIN),
            Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = objectives_data['north_star']
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = BODY_FONT_NAME
        subtitle_para.font.size = SUBTITLE_FONT_SIZE
        subtitle_para.font.color.rgb = BRAND_SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER


def create_objectives_slide(prs, objectives_data):
    """Create objectives slide displaying North star and key elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_BACKGROUND_COLOR
    
    # Add logo
    if LOGO_PATH:
        add_logo(slide, LOGO_PATH, LOGO_POSITION)
    
    # Slide title
    title_box = slide.shapes.add_textbox(
        SIDE_MARGIN,
        Inches(0.5),
        SLIDE_WIDTH - (2 * SIDE_MARGIN),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Objectives"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = TITLE_FONT_NAME
    title_para.font.size = HEADING_FONT_SIZE
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY_COLOR
    
    # North Star section
    y_pos = CONTENT_TOP_MARGIN
    
    if objectives_data.get('north_star'):
        north_star_box = slide.shapes.add_textbox(
            SIDE_MARGIN,
            y_pos,
            SLIDE_WIDTH - (2 * SIDE_MARGIN),
            Inches(1)
        )
        north_star_frame = north_star_box.text_frame
        north_star_frame.text = "North Star"
        north_star_para = north_star_frame.paragraphs[0]
        north_star_para.font.name = BODY_FONT_NAME
        north_star_para.font.size = Pt(24)
        north_star_para.font.bold = True
        north_star_para.font.color.rgb = BRAND_SECONDARY_COLOR
        
        y_pos += Inches(1.2)
        
        # North star content box
        if USE_SHAPES:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                SIDE_MARGIN,
                y_pos,
                SLIDE_WIDTH - (2 * SIDE_MARGIN),
                Inches(1.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = CONTENT_BOX_COLOR
            shape.line.color.rgb = BRAND_SECONDARY_COLOR
            shape.line.width = Pt(2)
            
            text_frame = shape.text_frame
            text_frame.text = objectives_data['north_star']
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            para = text_frame.paragraphs[0]
            para.font.name = BODY_FONT_NAME
            para.font.size = BODY_FONT_SIZE
            para.font.color.rgb = BRAND_TEXT_COLOR
        else:
            north_star_content = slide.shapes.add_textbox(
                SIDE_MARGIN,
                y_pos,
                SLIDE_WIDTH - (2 * SIDE_MARGIN),
                Inches(1.2)
            )
            north_star_content_frame = north_star_content.text_frame
            north_star_content_frame.text = objectives_data['north_star']
            north_star_content_para = north_star_content_frame.paragraphs[0]
            north_star_content_para.font.name = BODY_FONT_NAME
            north_star_content_para.font.size = BODY_FONT_SIZE
            north_star_content_para.font.color.rgb = BRAND_TEXT_COLOR
            north_star_content_frame.word_wrap = True
        
        y_pos += Inches(1.5)
    
    # Key Elements section
    if objectives_data.get('key_elements') and len(objectives_data['key_elements']) > 0:
        key_elements_title = slide.shapes.add_textbox(
            SIDE_MARGIN,
            y_pos,
            SLIDE_WIDTH - (2 * SIDE_MARGIN),
            Inches(0.6)
        )
        key_elements_title_frame = key_elements_title.text_frame
        key_elements_title_frame.text = "Key Elements"
        key_elements_title_para = key_elements_title_frame.paragraphs[0]
        key_elements_title_para.font.name = BODY_FONT_NAME
        key_elements_title_para.font.size = Pt(24)
        key_elements_title_para.font.bold = True
        key_elements_title_para.font.color.rgb = BRAND_SECONDARY_COLOR
        
        y_pos += Inches(0.8)
        
        # Key elements list
        elements_box = slide.shapes.add_textbox(
            SIDE_MARGIN + Inches(0.3),
            y_pos,
            SLIDE_WIDTH - (2 * SIDE_MARGIN) - Inches(0.3),
            SLIDE_HEIGHT - y_pos - BOTTOM_MARGIN
        )
        elements_frame = elements_box.text_frame
        elements_frame.word_wrap = True
        
        for i, element in enumerate(objectives_data['key_elements']):
            if i > 0:
                elements_frame.add_paragraph()
            para = elements_frame.paragraphs[i]
            para.text = f"• {str(element)}"
            para.font.name = BODY_FONT_NAME
            para.font.size = BODY_FONT_SIZE
            para.font.color.rgb = BRAND_TEXT_COLOR
            para.space_after = Pt(8)
            para.level = 0


def create_roadmap_slides(prs, roadmap_df):
    """Create roadmap slides grouped by timeline/phase."""
    if roadmap_df.empty:
        return
    
    # Group by Timeline
    grouped = roadmap_df.groupby('Timeline')
    
    for timeline, group in grouped:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BRAND_BACKGROUND_COLOR
        
        # Add logo
        if LOGO_PATH:
            add_logo(slide, LOGO_PATH, LOGO_POSITION)
        
        # Slide title (Timeline)
        title_box = slide.shapes.add_textbox(
            SIDE_MARGIN,
            Inches(0.5),
            SLIDE_WIDTH - (2 * SIDE_MARGIN),
            Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Roadmap: {str(timeline)}"
        title_para = title_frame.paragraphs[0]
        title_para.font.name = TITLE_FONT_NAME
        title_para.font.size = HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = BRAND_PRIMARY_COLOR
        
        # Group by Phase within this timeline
        phase_groups = group.groupby('Phase') if 'Phase' in group.columns else [(None, group)]
        
        y_pos = CONTENT_TOP_MARGIN
        max_width = (SLIDE_WIDTH - (2 * SIDE_MARGIN)) / len(phase_groups) if len(phase_groups) > 1 else (SLIDE_WIDTH - (2 * SIDE_MARGIN))
        box_width = max_width - Inches(0.2)
        
        for phase_idx, (phase, phase_data) in enumerate(phase_groups):
            x_pos = SIDE_MARGIN + (phase_idx * max_width) + Inches(0.1)
            
            # Phase header
            if phase and str(phase).strip():
                phase_title = slide.shapes.add_textbox(
                    x_pos,
                    y_pos,
                    box_width,
                    Inches(0.6)
                )
                phase_title_frame = phase_title.text_frame
                phase_title_frame.text = str(phase)
                phase_title_para = phase_title_frame.paragraphs[0]
                phase_title_para.font.name = BODY_FONT_NAME
                phase_title_para.font.size = Pt(22)
                phase_title_para.font.bold = True
                phase_title_para.font.color.rgb = BRAND_SECONDARY_COLOR
                y_pos_phase = y_pos + Inches(0.7)
            else:
                y_pos_phase = y_pos
            
            # Workpackages for this phase
            workpackages = phase_data['Workpackage'].dropna().tolist()
            
            if workpackages:
                # Create content box for workpackages
                content_height = min(Inches(4), len(workpackages) * Inches(0.8) + Inches(0.3))
                
                if USE_SHAPES:
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        x_pos,
                        y_pos_phase,
                        box_width,
                        content_height
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = CONTENT_BOX_COLOR
                    shape.line.color.rgb = BRAND_ACCENT_COLOR
                    shape.line.width = Pt(1.5)
                    
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True
                    text_frame.margin_left = Inches(0.15)
                    text_frame.margin_right = Inches(0.15)
                    text_frame.margin_top = Inches(0.15)
                    text_frame.margin_bottom = Inches(0.15)
                    
                    for i, wp in enumerate(workpackages):
                        if i > 0:
                            text_frame.add_paragraph()
                        para = text_frame.paragraphs[i]
                        para.text = f"• {str(wp)}"
                        para.font.name = BODY_FONT_NAME
                        para.font.size = Pt(14)
                        para.font.color.rgb = BRAND_TEXT_COLOR
                        para.space_after = Pt(6)
                else:
                    wp_box = slide.shapes.add_textbox(
                        x_pos,
                        y_pos_phase,
                        box_width,
                        content_height
                    )
                    wp_frame = wp_box.text_frame
                    wp_frame.word_wrap = True
                    
                    for i, wp in enumerate(workpackages):
                        if i > 0:
                            wp_frame.add_paragraph()
                        para = wp_frame.paragraphs[i]
                        para.text = f"• {str(wp)}"
                        para.font.name = BODY_FONT_NAME
                        para.font.size = Pt(14)
                        para.font.color.rgb = BRAND_TEXT_COLOR
                        para.space_after = Pt(6)


def main():
    """Main function to orchestrate PowerPoint generation."""
    parser = argparse.ArgumentParser(
        description='Generate PowerPoint presentation from Excel roadmap file'
    )
    parser.add_argument(
        'excel_file',
        type=str,
        help='Path to Excel file containing Objectives and Roadmap sheets'
    )
    parser.add_argument(
        '-o', '--output',
        type=str,
        default=None,
        help='Output PowerPoint file path (default: same as Excel file with .pptx extension)'
    )
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.exists(args.excel_file):
        print(f"Error: Excel file '{args.excel_file}' not found.")
        return
    
    # Determine output file path
    if args.output:
        output_path = args.output
    else:
        base_name = os.path.splitext(args.excel_file)[0]
        output_path = f"{base_name}.pptx"
    
    print(f"Reading Excel file: {args.excel_file}")
    
    # Read Excel data
    objectives_data = read_objectives(args.excel_file)
    roadmap_df = read_roadmap(args.excel_file)
    
    print(f"Found {len(roadmap_df)} roadmap entries")
    print(f"North Star: {objectives_data.get('north_star', 'Not found')}")
    print(f"Key Elements: {len(objectives_data.get('key_elements', []))} items")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Generate slides
    print("Generating slides...")
    create_title_slide(prs, objectives_data)
    create_objectives_slide(prs, objectives_data)
    create_roadmap_slides(prs, roadmap_df)
    
    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
